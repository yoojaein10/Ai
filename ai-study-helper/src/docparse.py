"""문서 파서: PDF / Word / PPT / Excel / TXT / MD → 위치 정보가 붙은 텍스트 세그먼트.

세그먼트마다 위치(페이지·슬라이드·시트)를 보존하는 이유:
이후 RAG Q&A에서 "이 답의 근거는 p.3"처럼 출처를 표시해야 하기 때문(계획서 2장 6번).

읽을 수 없는 문서(스캔 PDF 등)는 is_readable=False 로 표시만 해둔다.
실제 OCR 폴백(로컬 OCR → Gemini 비전 2단)은 비전 래퍼가 생기는 3단계 이후 구현.
"""

from dataclasses import dataclass, field
from pathlib import Path

from . import config

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".md"}


@dataclass
class DocSegment:
    location: str  # 예: "p.3", "슬라이드 2", "시트:매출"
    text: str


@dataclass
class ParsedDocument:
    path: str
    doc_type: str  # pdf | docx | pptx | xlsx | text
    segments: list[DocSegment] = field(default_factory=list)
    is_readable: bool = True  # False면 언리더블 → OCR 폴백 대상(3단계 이후)

    @property
    def full_text(self) -> str:
        return "\n\n".join(f"[{s.location}]\n{s.text}" for s in self.segments if s.text.strip())

    @property
    def char_count(self) -> int:
        return sum(len(s.text) for s in self.segments)


def parse_document(path: str | Path, ocr_fallback: bool = True) -> ParsedDocument:
    """확장자에 맞는 파서로 문서를 읽고, 텍스트가 거의 없으면 언리더블로 표시한다.

    ocr_fallback=True면 언리더블 PDF에 2단 OCR 폴백을 시도한다:
    1차 로컬 OCR(Tesseract, 설치된 경우) → 결과 부실 시 2차 Gemini 비전.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"문서를 찾을 수 없습니다: {path}")
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(f"지원하지 않는 문서 형식입니다: {ext} (지원: {sorted(SUPPORTED_EXTS)})")

    parser = {
        ".pdf": _parse_pdf,
        ".docx": _parse_docx,
        ".pptx": _parse_pptx,
        ".xlsx": _parse_xlsx,
        ".txt": _parse_text,
        ".md": _parse_text,
    }[ext]
    doc = parser(path)
    doc.is_readable = _looks_readable(doc)

    if not doc.is_readable and ocr_fallback and ext == ".pdf":
        # 스캔 PDF — 페이지를 이미지로 렌더링해 OCR 경로로 재추출
        doc = ocr_pdf(path)
        doc.is_readable = _looks_readable(doc)
    return doc


def detect_language(text: str) -> str:
    """간단한 언어 감지: 한글/라틴 문자 비율. 반환: "ko" | "en" | "unknown".

    라이브러리 없이 휴리스틱으로 충분한 이유: 분기 목적이 "한국어냐 아니냐"
    둘뿐이고, 학습 자료는 혼용돼도 한글이 조금만 있으면 한국어 자료다."""
    sample = text[:5000]
    hangul = sum(1 for c in sample if "가" <= c <= "힣")
    latin = sum(1 for c in sample if c.isascii() and c.isalpha())
    if hangul + latin < 20:
        return "unknown"
    if hangul >= (hangul + latin) * 0.15:  # 한글이 15% 이상이면 한국어 자료로 본다
        return "ko"
    return "en"


# 논문 구조 신호 — 이 중 3개 이상 발견되면 논문으로 판별
_PAPER_MARKERS = [
    "abstract", "introduction", "related work", "methodology", "method",
    "experiment", "results", "conclusion", "references", "초록", "서론", "결론", "참고문헌",
]


def looks_like_paper(text: str) -> bool:
    """논문 판별: 섹션 표제 신호가 3개 이상이면 논문 구조 요약 경로를 활성화한다."""
    lowered = text[:20000].lower()
    return sum(1 for m in _PAPER_MARKERS if m in lowered) >= 3


def _looks_readable(doc: ParsedDocument) -> bool:
    """언리더블 판별: 페이지(세그먼트)당 평균 글자 수가 임계치 미만이면 스캔본으로 간주."""
    if not doc.segments:
        return False
    return doc.char_count / len(doc.segments) >= config.MIN_CHARS_PER_PAGE


# ---------- 삽입 도표·수식 비전 해석 (계획서 10-1: 논문 멀티모달) ----------

FIGURE_PROMPT = """이 이미지는 문서(논문·강의자료)에 삽입된 그림입니다.
도표·그래프면 축·수치·경향과 의미를, 수식이면 수식과 뜻을, 다이어그램이면 구조를
한국어로 3~6문장 설명하세요. 장식용 이미지면 "장식 이미지"라고만 답하세요."""

MAX_FIGURES_PER_DOC = 8  # 비전 호출 상한 — 무료 티어 보호
# 도표는 플랫 컬러가 많아 PNG 압축이 잘 된다(실측: 625x417 차트가 6KB) —
# 바이트 임계값은 낮게 두고, 아이콘 제거는 치수 조건(MIN_FIGURE_DIM)이 담당한다.
MIN_FIGURE_BYTES = 3 * 1024
MIN_FIGURE_DIM = 120


def _extract_pdf_figures(path: Path, out_dir: Path) -> list[tuple[int, Path]]:
    """PDF에 삽입된 의미 있는 크기의 이미지들을 (페이지 번호, 파일) 목록으로 추출."""
    import fitz

    out_dir.mkdir(parents=True, exist_ok=True)
    figures: list[tuple[int, Path]] = []
    with fitz.open(path) as pdf:
        for page_no, page in enumerate(pdf, start=1):
            for img_index, img in enumerate(page.get_images(full=True)):
                if len(figures) >= MAX_FIGURES_PER_DOC:
                    return figures
                pix = fitz.Pixmap(pdf, img[0])
                if pix.width < MIN_FIGURE_DIM or pix.height < MIN_FIGURE_DIM:
                    continue
                if pix.n > 4:  # CMYK 등 → RGB 변환
                    pix = fitz.Pixmap(fitz.csRGB, pix)
                fig_path = out_dir / f"fig_p{page_no}_{img_index}.png"
                pix.save(str(fig_path))
                if fig_path.stat().st_size < MIN_FIGURE_BYTES:
                    fig_path.unlink()
                    continue
                figures.append((page_no, fig_path))
    return figures


def analyze_pdf_figures(path: Path, doc: ParsedDocument) -> int:
    """삽입 도표를 Gemini 비전으로 해석해 해당 페이지 세그먼트에 병합한다.

    반환: 해석된 그림 수. 텍스트만으로는 잃어버리는 논문의 도표·수식 정보를
    기존 비전 래퍼 재활용으로 살리는 단계 (계획서 10-1)."""
    import tempfile

    by_location = {s.location: s for s in doc.segments}
    analyzed = 0
    with tempfile.TemporaryDirectory(prefix="figs_") as tmp:
        figures = _extract_pdf_figures(path, Path(tmp))
        if not figures:
            return 0
        from .llm import get_client

        client = get_client()
        for page_no, fig_path in figures:
            try:
                desc = client.generate(FIGURE_PROMPT, images=[fig_path], fast=True).strip()
            except Exception as e:
                print(f"  그림 해석 실패(p.{page_no}): {str(e)[:80]}")
                continue
            if "장식 이미지" in desc[:20]:
                continue
            seg = by_location.get(f"p.{page_no}")
            if seg is not None:
                seg.text = f"{seg.text}\n\n[그림 해석 (p.{page_no})]\n{desc}"
                analyzed += 1
    if analyzed:
        print(f"삽입 도표 {analyzed}개 비전 해석 완료")
    return analyzed


# ---------- 2단 OCR 폴백 (계획서 4-1: 언리더블 문서) ----------

OCR_PROMPT = """이 이미지는 문서 페이지의 스캔본입니다. 보이는 텍스트를 빠짐없이 추출하세요.
- 원문 그대로, 줄바꿈·목록 구조 유지
- 표는 마크다운 표로
- 수식은 텍스트 표기로 (예: E = mc^2)
- 그림·도장·서명은 무시
- 추출한 텍스트만 출력 (설명 금지)"""


def _render_pdf_pages(path: Path, out_dir: Path, dpi: int = 150) -> list[Path]:
    """PDF 페이지들을 OCR용 이미지로 렌더링."""
    import fitz

    out_dir.mkdir(parents=True, exist_ok=True)
    images: list[Path] = []
    with fitz.open(path) as pdf:
        for i, page in enumerate(pdf, start=1):
            pix = page.get_pixmap(dpi=dpi)
            img_path = out_dir / f"page_{i:03d}.png"
            pix.save(str(img_path))
            images.append(img_path)
    return images


def _ocr_local(image_path: Path) -> str | None:
    """1차: 로컬 OCR (Tesseract). 미설치 환경이면 None — 2차로 넘어간다."""
    try:
        import pytesseract
        from PIL import Image

        return pytesseract.image_to_string(Image.open(image_path), lang="kor+eng")
    except Exception:
        return None


def ocr_pdf(path: Path) -> ParsedDocument:
    """스캔 PDF → 페이지 이미지 → 2단 OCR (1차 로컬 → 부실하면 2차 Gemini 비전).

    2단 구조인 이유: 로컬 OCR은 비용이 없지만 복잡한 표·수식·저화질에 약하다.
    페이지별로 판정해 필요한 페이지만 Gemini를 호출한다(호출량 최소화)."""
    import tempfile

    doc = ParsedDocument(path=str(path), doc_type="pdf")
    with tempfile.TemporaryDirectory(prefix="ocr_") as tmp:
        images = _render_pdf_pages(path, Path(tmp))
        gemini_used = 0
        for i, img in enumerate(images, start=1):
            text = _ocr_local(img)
            if not text or len(text.strip()) < config.MIN_CHARS_PER_PAGE:
                from .llm import get_client  # 지연 import — OCR 안 쓰는 경로에선 로드 안 됨

                try:
                    text = get_client().generate(OCR_PROMPT, images=[img], fast=True)
                    gemini_used += 1
                except Exception as e:
                    # 빈 페이지(추출할 글자 없음) 등 — 한 페이지 실패가 문서 전체를 막으면 안 된다
                    print(f"  p.{i} OCR 실패, 빈 페이지로 처리: {str(e)[:80]}")
                    text = ""
            doc.segments.append(DocSegment(location=f"p.{i}", text=(text or "").strip()))
        if gemini_used:
            print(f"OCR 폴백: {len(images)}페이지 중 {gemini_used}페이지에 Gemini 비전 사용")
    return doc


def _parse_pdf(path: Path) -> ParsedDocument:
    import fitz  # PyMuPDF

    doc = ParsedDocument(path=str(path), doc_type="pdf")
    with fitz.open(path) as pdf:
        for i, page in enumerate(pdf, start=1):
            doc.segments.append(DocSegment(location=f"p.{i}", text=page.get_text().strip()))
    return doc


def _parse_docx(path: Path) -> ParsedDocument:
    import docx

    d = docx.Document(str(path))
    doc = ParsedDocument(path=str(path), doc_type="docx")
    parts: list[str] = [p.text for p in d.paragraphs if p.text.strip()]
    # 표는 행을 탭 구분 텍스트로 펼친다 — LLM이 표 구조를 인식할 수 있는 형태
    for table in d.tables:
        rows = ["\t".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        parts.append("\n".join(rows))
    # Word는 페이지 경계를 파일에 저장하지 않으므로 문서 전체를 한 세그먼트로 다룬다
    doc.segments.append(DocSegment(location="본문", text="\n".join(parts)))
    return doc


def _parse_pptx(path: Path) -> ParsedDocument:
    from pptx import Presentation

    prs = Presentation(str(path))
    doc = ParsedDocument(path=str(path), doc_type="pptx")
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if shape.has_table:
                rows = ["\t".join(c.text.strip() for c in row.cells) for row in shape.table.rows]
                parts.append("\n".join(rows))
        # 발표자 노트에는 슬라이드에 없는 설명이 담기는 경우가 많다
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append(f"(발표자 노트) {slide.notes_slide.notes_text_frame.text.strip()}")
        doc.segments.append(DocSegment(location=f"슬라이드 {i}", text="\n".join(parts)))
    return doc


def _parse_xlsx(path: Path) -> ParsedDocument:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)  # data_only: 수식 대신 계산값
    doc = ParsedDocument(path=str(path), doc_type="xlsx")
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            if any(v is not None for v in row):
                rows.append("\t".join("" if v is None else str(v) for v in row))
        doc.segments.append(DocSegment(location=f"시트:{sheet.title}", text="\n".join(rows)))
    wb.close()
    return doc


def _parse_text(path: Path) -> ParsedDocument:
    doc = ParsedDocument(path=str(path), doc_type="text")
    doc.segments.append(DocSegment(location="본문", text=path.read_text(encoding="utf-8", errors="replace")))
    return doc
