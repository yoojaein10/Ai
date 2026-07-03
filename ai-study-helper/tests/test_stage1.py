"""1단계 테스트: 입력 유형 감지, 문서 파싱, job_id 캐싱.

문서 픽스처(PDF/Word/PPT/Excel)는 파서와 같은 라이브러리로 즉석 생성한다 —
바이너리 픽스처를 저장소에 커밋하지 않아도 되고, 내용을 코드에서 바로 볼 수 있다.
"""

import pytest

from src.docparse import parse_document
from src.ingest import detect_source_type, make_job_id


# --- 입력 유형 감지 ---------------------------------------------------------

@pytest.mark.parametrize(
    "source, expected",
    [
        ("D:/videos/lecture.mp4", "video"),
        ("lecture.MKV", "video"),
        ("audio.mp3", "audio"),
        ("https://www.youtube.com/watch?v=abc123", "youtube"),
        ("https://youtu.be/abc123", "youtube"),
        ("notes.pdf", "document"),
        ("slides.pptx", "document"),
        ("data.xlsx", "document"),
        ("memo.md", "document"),
    ],
)
def test_detect_source_type(source, expected):
    assert detect_source_type(source) == expected


def test_detect_unsupported_raises():
    with pytest.raises(ValueError):
        detect_source_type("archive.zip")


# --- job_id: 같은 입력 → 같은 id (캐싱의 전제) ------------------------------

def test_job_id_stable_for_same_file(tmp_path):
    f = tmp_path / "a.txt"
    f.write_text("동일한 내용", encoding="utf-8")
    assert make_job_id(str(f), "document") == make_job_id(str(f), "document")


def test_job_id_differs_for_different_content(tmp_path):
    f1 = tmp_path / "a.txt"
    f2 = tmp_path / "b.txt"
    f1.write_text("내용 A", encoding="utf-8")
    f2.write_text("내용 B", encoding="utf-8")
    assert make_job_id(str(f1), "document") != make_job_id(str(f2), "document")


def test_job_id_for_youtube_is_url_based():
    url = "https://www.youtube.com/watch?v=abc123"
    assert make_job_id(url, "youtube") == make_job_id(url, "youtube")


# --- 문서 파서 --------------------------------------------------------------

def test_parse_txt(tmp_path):
    f = tmp_path / "memo.txt"
    f.write_text("머신러닝은 데이터로부터 패턴을 학습하는 기술이다.\n" * 5, encoding="utf-8")
    doc = parse_document(f)
    assert doc.doc_type == "text"
    assert doc.is_readable
    assert "머신러닝" in doc.full_text


def test_parse_pdf(tmp_path):
    import fitz

    f = tmp_path / "sample.pdf"
    pdf = fitz.open()
    for i in range(2):
        page = pdf.new_page()
        page.insert_text((72, 72), f"Deep learning lecture page {i + 1}. " * 5)
    pdf.save(str(f))
    pdf.close()

    doc = parse_document(f)
    assert doc.doc_type == "pdf"
    assert len(doc.segments) == 2
    assert doc.segments[0].location == "p.1"
    assert "Deep learning" in doc.full_text
    assert doc.is_readable


def test_parse_pdf_unreadable_detection(tmp_path):
    """빈 페이지(스캔본 시뮬레이션)는 언리더블로 판별되어야 한다."""
    import fitz

    f = tmp_path / "scanned.pdf"
    pdf = fitz.open()
    pdf.new_page()
    pdf.new_page()
    pdf.save(str(f))
    pdf.close()

    doc = parse_document(f, ocr_fallback=False)  # 판별만 검증 (폴백은 11단계 테스트)
    assert not doc.is_readable


def test_parse_docx_with_table(tmp_path):
    import docx

    f = tmp_path / "sample.docx"
    d = docx.Document()
    d.add_paragraph("신경망의 기초 개념 정리 문서입니다.")
    table = d.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "용어"
    table.cell(0, 1).text = "설명"
    table.cell(1, 0).text = "역전파"
    table.cell(1, 1).text = "오차를 뒤로 전파해 가중치를 갱신"
    d.save(str(f))

    doc = parse_document(f)
    assert doc.doc_type == "docx"
    assert "신경망" in doc.full_text
    assert "역전파\t오차를 뒤로 전파해 가중치를 갱신" in doc.full_text


def test_parse_pptx_with_notes(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    f = tmp_path / "sample.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 제목만 있는 레이아웃
    slide.shapes.title.text = "1장: 지도학습이란"
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    box.text_frame.text = "정답(레이블)이 있는 데이터로 모델을 학습시킨다."
    slide.notes_slide.notes_text_frame.text = "여기서 회귀와 분류의 차이를 설명할 것."
    prs.save(str(f))

    doc = parse_document(f)
    assert doc.doc_type == "pptx"
    assert doc.segments[0].location == "슬라이드 1"
    assert "지도학습" in doc.full_text
    assert "(발표자 노트) 여기서 회귀와 분류의 차이를 설명할 것." in doc.full_text


def test_parse_xlsx(tmp_path):
    from openpyxl import Workbook

    f = tmp_path / "sample.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "성적"
    ws.append(["이름", "점수"])
    ws.append(["김철수", 85])
    ws.append(["이영희", 92])
    wb.save(str(f))

    doc = parse_document(f)
    assert doc.doc_type == "xlsx"
    assert doc.segments[0].location == "시트:성적"
    assert "김철수\t85" in doc.full_text


def test_parse_unsupported_ext_raises(tmp_path):
    f = tmp_path / "archive.zip"
    f.write_bytes(b"PK")
    with pytest.raises(ValueError):
        parse_document(f)


def test_parse_missing_file_raises():
    with pytest.raises(FileNotFoundError):
        parse_document("없는파일.pdf")
